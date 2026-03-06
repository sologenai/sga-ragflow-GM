import os
import google.generativeai as genai
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import time
import base64

# --- CONFIGURATION ---
BASE_PATH = r"e:\sga-ragflow-GM\ragflow_training_ppt"
IMAGE_PATH = os.path.join(BASE_PATH, "images")
PPT_FILE = os.path.join(BASE_PATH, "RAGFlow_Training_XiamenITG.pptx")

if not os.path.exists(IMAGE_PATH):
    os.makedirs(IMAGE_PATH)

# OpenRouter Configuration
OPENROUTER_API_KEY = "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"
openai_client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=OPENROUTER_API_KEY
)

# --- STYLE CONFIGURATION (Professional Template) ---
STYLE_CORE = (
    "风格：高端商务科技风。主题：RAGFlow 企业培训。背景：深海蓝 (#0A192F) 搭配隐约的几何网格线条。光效：3D 玻璃拟态元素，柔和的紫色和电光青色光。渲染：干净、极简、电影级 16:9。"
)

# --- SLIDE DATA DEFINITION (2026 厦门国贸定制版 - 全干货) ---
slides_data = [
    {
        "id": 1, 
        "title": "封面：厦门国贸 RAGFlow 开发管理员培训", 
        "subtitle": "构建企业级供应链知识大脑\n2026年3月",
        "content": [],
        "img_prompt": "专业技术架构图。画面上方：3D发光白色中文字“厦门国贸 RAGFlow 开发管理员培训”。画面下方：3D发光白色中文字“2026年3月”。画面中央：RAGFlow 系统架构框图，包含以下模块：文档输入 → DeepDoc解析 → 切片 → Embedding → 向量库 → 检索 → LLM → 答案。每个模块用中文标注。线条用青色和紫色表示数据流。"
    },
    {
        "id": 2, 
        "title": "课程目标：从文档到智能问答", 
        "content": [
            "• 目标1：掌握 DeepDoc 深度文档解析的12种解析模板",
            "• 目标2：理解混合检索+Rerank的完整链路",
            "• 目标3：能够配置知识图谱进行多跳推理",
            "• 目标4：通过 Agent 工作流实现业务自动化"
        ],
        "img_prompt": "技术流程图，包含四个阶段，每个阶段用中文标注：第一阶段：文档解析（DeepDoc），配图标文档和放大镜；第二阶段：混合检索（BM25+向量），配图标搜索；第三阶段：知识图谱（GraphRAG），配图标网络节点；第四阶段：智能代理（Agent），配图标机器人。每个阶段下方写中文：掌握解析模板、理解检索链路、配置知识图谱、实现业务自动化。"
    },
    {
        "id": 3, 
        "title": "RAGFlow 核心架构", 
        "content": [
            "• 文档层：支持 PDF/Word/Excel/PPT/扫描件/网页等12+格式",
            "• 解析层：DeepDoc + MinerU + Docling，布局识别+表格还原",
            "• 检索层：向量检索+关键词检索+图谱检索，多路召回",
            "• 应用层：对话 + Agent + API + 多终端集成"
        ],
        "img_prompt": "完整的系统架构图，分四层展示。最上层写中文：应用层，包含：对话、Agent、API、多终端。第二层写中文：检索层，包含：向量检索、关键词检索、图谱检索。第三层写中文：解析层，包含：DeepDoc、MinerU、Docling、布局识别、表格还原。最下层写中文：文档层，包含：PDF、Word、Excel、PPT、扫描件、网页。用青色和紫色箭头连接各层，形成数据流。"
    },
    {
        "id": 4, 
        "title": "什么是 RAG：检索增强生成", 
        "content": [
            "• RAG = 检索(Retrieval) + 生成(Generation)",
            "• 解决大模型两个核心问题：幻觉 + 知识过期",
            "• 工作流：用户问题 → 检索相关文档 → 作为上下文 → LLM生成答案",
            "• 核心优势：答案有据可查，可溯源原文"
        ],
        "img_prompt": "RAG工作流程图，从左到右：第一步：用户问题，配图标问号，写中文“用户问题”；第二步：文档检索，配图标搜索，写中文“检索相关文档”；第三步：上下文拼接，配图标拼图，写中文“作为上下文”；第四步：LLM生成，配图标大脑，写中文“生成答案”；第五步：输出答案，配图标对勾，写中文“有据可查的答案”。上方大标题写中文：RAG = 检索 + 生成。"
    },
    {
        "id": 5, 
        "title": "DeepDoc 核心：12种解析模板详解", 
        "content": [
            "• Naive：通用模式，智能分段（规章制度）",
            "• Table：表格专家，还原表头和数据对应关系",
            "• Laws：法律法规，自动识别章-节-条-款",
            "• Manual：产品手册，识别步骤、列表、表格",
            "• Paper：学术论文，识别标题、摘要、参考文献"
        ],
        "img_prompt": "解析模板对比图，5个并排的卡片。第一个卡片：Naive，写中文“通用模式-规章制度”；第二个卡片：Table，写中文“表格专家-财务报表”；第三个卡片：Laws，写中文“法律法规-合同条款”；第四个卡片：Manual，写中文“产品手册-操作指南”；第五个卡片：Paper，写中文“学术论文-技术文档”。每个卡片配相应图标。上方大标题写中文：DeepDoc 12种解析模板。"
    },
    {
        "id": 6, 
        "title": "Table 模板实战：财务报表解析", 
        "content": [
            "• 挑战：跨页表格、合并单元格、无框线表格",
            "• 配置：chunk_token_num=300, table_min_rows=2",
            "• 效果：PDF表格 → Markdown/HTML结构，保留行列关系",
            "• 验证：确保表头与数据对齐，公式和数字准确"
        ],
        "img_prompt": "表格解析流程图，左侧：原始PDF表格（有合并单元格和跨页），写中文“原始PDF表格”；中间：解析过程，写中文“DeepDoc Table模板解析”，配图标齿轮；右侧：结构化输出，写中文“Markdown/HTML结构”，显示整齐的表格。下方写配置参数：chunk_token_num=300, table_min_rows=2。"
    },
    {
        "id": 7, 
        "title": "Laws 模板：合同与法律法规", 
        "content": [
            "• 场景：采购合同、ISO标准、监管文件",
            "• 能力：自动提取“章-节-条-款”层级结构",
            "• 检索优势：保留父级标题，明确上下文",
            "• 示例：第一章 总则 → 第一条 定义 → （一）供应链"
        ],
        "img_prompt": "层级结构图，树状结构。根节点写中文“合同文档”，第一层：第一章 总则、第二章 采购、第三章 付款；第二层：第一条 定义、第二条 范围、第三条 供应商；第三层：（一）供应链、（二）货物、（三）服务。用连接线连接，形成清晰的层级。配说明文字：Laws模板自动识别章-节-条-款。"
    },
    {
        "id": 8, 
        "title": "切片参数配置：Chunk Size & Overlap", 
        "content": [
            "• Chunk Size：推荐 300-512 Tokens（中文约200-350字）",
            "• Chunk Overlap：推荐 10%-15%，约50-100 Tokens",
            "• 过小：语义破碎，丢失语境",
            "• 过大：包含噪音，检索不准"
        ],
        "img_prompt": "参数配置示意图，上方：Chunk Size滑块，从左到右标注：300、512、1024，箭头指向512，写中文“推荐：300-512 Tokens”。下方：Overlap示意图，两个矩形部分重叠，重叠区域高亮，标注10%-15%，写中文“Overlap：50-100 Tokens”。配说明：语义胶水，避免关键词被切断。"
    },
    {
        "id": 9, 
        "title": "布局识别：智能去噪", 
        "content": [
            "• 自动识别并剔除：页眉、页脚、页码、侧边栏",
            "• 识别非正文内容：水印、广告、导航链接",
            "• 配置项：layout_recognize: true（DeepDoc）",
            "• 效果：确保喂给大模型的是纯正文内容"
        ],
        "img_prompt": "文档去噪前后对比图。左侧：原始文档，标注页眉、页脚、水印、侧边栏，用红色叉号标记，写中文“噪音内容”。右侧：清理后的文档，只保留正文，用绿色对勾标记，写中文“纯净正文”。中间箭头写中文“DeepDoc布局识别”。配置参数：layout_recognize: true。"
    },
    {
        "id": 10, 
        "title": "多模态解析：OCR与图表理解", 
        "content": [
            "• OCR：PaddleOCR/VL，支持扫描件文字提取",
            "• 图表理解：多模态模型识别流程图、架构图",
            "• 配置：image2text_model=gpt-4-vision",
            "• 适用：扫描合同、发票、技术图纸、产品图片"
        ],
        "img_prompt": "多模态解析流程图，左侧：输入（扫描合同、发票、技术图纸）配相应图标。中间：OCR模块，写中文“PaddleOCR文字提取”；图表理解模块，写中文“多模态模型识别”。右侧：输出（结构化文字、图表描述）。配置参数：image2text_model=gpt-4-vision。"
    },
    {
        "id": 11, 
        "title": "混合检索：关键词 + 向量 + 重排", 
        "content": [
            "• 关键词检索（BM25）：精确匹配，专有名词效果好",
            "• 向量检索（Embedding）：语义匹配，模糊意图效果好",
            "• Rerank重排：对Top30候选二次排序，提升Top5准确率",
            "• 权重配置：关键词0.3，向量0.7，相似度阈值0.2-0.5"
        ],
        "img_prompt": "混合检索架构图，三个并行输入流：第一个流：关键词检索，写中文“BM25精确匹配”，权重0.3；第二个流：向量检索，写中文“Embedding语义匹配”，权重0.7；第三个流：图谱检索，写中文“GraphRAG关系推理”。三个流汇聚到Rerank模块，写中文“重排模型bge-reranker-large Top30→Top10”。最终输出：精准答案。"
    },
    {
        "id": 12, 
        "title": "Embedding 模型选型：中文场景", 
        "content": [
            "• 首选：Qwen3-Embedding-VL（阿里开源，多模态适配，支持图像和文本）",
            "• 备选：Qwen3-Embedding（纯文本，更快）",
            "• 备选：BAAI/bge-large-zh-v1.5（中文效果稳定）",
            "• 重要：同一知识库必须使用相同Embedding模型，升级需重新向量化"
        ],
        "img_prompt": "Embedding模型对比表，三个卡片并排。第一个卡片：Qwen3-Embedding-VL，写中文\"推荐首选-多模态适配-支持图像和文本\"，星级★★★★★；第二个卡片：Qwen3-Embedding，写中文\"纯文本-速度更快\"，星级★★★★☆；第三个卡片：BGE-Large-zh-v1.5，写中文\"中文效果稳定\"，星级★★★☆☆。底部红色警告框写中文：同一知识库必须使用相同Embedding模型，升级需重新向量化！"
    },
    {
        "id": 13, 
        "title": "检索测试与调优", 
        "content": [
            "• 测试页面：知识库→Test标签，输入测试问题",
            "• 评估指标：Top-K召回率、相似度分数、排序质量",
            "• 调优方向：调整Top-K、相似度阈值、权重配置",
            "• 目标：Top-5召回率>90%，最相关结果排第一"
        ],
        "img_prompt": "检索测试界面示意图，左侧：测试问题输入框，写中文“测试问题：2025年采购政策是什么？”。中间：检索结果列表，显示Top-10结果，每个结果有相似度分数。右侧：评估指标面板，显示Top-5召回率92%、MRR 0.85、NDCG 0.88。底部写中文：目标 - Top-5召回率>90%。"
    },
    {
        "id": 14, 
        "title": "知识图谱 GraphRAG：多跳推理", 
        "content": [
            "• 实体类型：组织、人员、地理位置、事件、产品",
            "• 关系类型：供应、采购、隶属、存储、合作",
            "• 构建方法：Light（轻量，省token）/ General（通用，更准确）",
            "• 场景：供应商风险传导、股权穿透分析、供应链全景"
        ],
        "img_prompt": "供应链知识图谱示意图，节点包括：厦门国贸（蓝色）、供应商A（绿色）、供应商B（绿色）、客户C（橙色）、仓库D（紫色）。边标注：供应、采购、存储。配说明文字：实体类型-组织/人员/地理位置；关系类型-供应/采购/隶属。下方写中文：GraphRAG多跳推理：厦门国贸的供应商A的客户是谁？"
    },
    {
        "id": 15, 
        "title": "GraphRAG 配置实战", 
        "content": [
            "• 启用：创建知识库→高级设置→Enable GraphRAG",
            "• 实体类型：组织、人员、地理位置、产品、合同",
            "• 方法：Light（推荐，省token）",
            "• 实体消歧：启用，解决“国贸”与“厦门国贸”同一实体"
        ],
        "img_prompt": "GraphRAG配置界面截图风格，左侧：配置表单，包含复选框Enable GraphRAG（已勾选）、实体类型输入框（组织,人员,地理位置,产品,合同）、方法选择（Light/General，选Light）、实体消歧（已勾选）。右侧：图谱预览，显示实体节点和关系连线。底部写中文：构建图谱需要时间和token，建议先在小数据集测试。"
    },
    {
        "id": 16, 
        "title": "原文溯源：可信 AI", 
        "content": [
            "• 拒绝黑盒：每一次回答都有据可查",
            "• 引用快照：答案旁显示相关原文片段",
            "• 跳转链接：点击引用直接跳转至原文具体位置",
            "• 反馈闭环：用户可标注答案正确性，持续优化"
        ],
        "img_prompt": "溯源界面示意图，上方：问题“2025年付款政策变更内容是什么？”。中间：LLM答案，用蓝色高亮标注引用标记[1][2]。下方：引用快照区域，显示两个片段：[1] 2025年3月发布的《付款政策变更通知》原文片段；[2] 2025年4月财务部邮件原文片段。配说明文字：点击引用可跳转原文。"
    },
    {
        "id": 17, 
        "title": "Agent 工作流：从问答到自动化", 
        "content": [
            "• Agent = LLM + 知识库 + 工具 + 记忆",
            "• 核心组件：Begin→Retrieval→Generate→Answer",
            "• 高级组件：Categorize（分类）、Switch（分支）、ExeSQL",
            "• 模板：Simple RAG、Deep RAG、Self RAG、Agentic RAG"
        ],
        "img_prompt": "Agent工作流画布图，节点用框和箭头连接。开始节点：Begin；然后是Retrieval（检索知识库）；然后是Generate（LLM生成）；然后是Answer（输出答案）。旁边列出高级组件：Categorize分类、Switch分支、ExeSQL执行SQL、ExePython执行代码、HTTP请求。下方写中文：预置模板：Simple RAG、Deep RAG、Self RAG、Agentic RAG。"
    },
    {
        "id": 18, 
        "title": "System Prompt 工程：人设与约束", 
        "content": [
            "• 人设定义：你是厦门国贸资深供应链风控专家",
            "• 约束条件：只回答知识库内内容，不知道说不知道",
            "• 输出格式：请以表格形式输出风险点对比",
            "• 参考：框架级Prompt块 - task_analysis、reflect"
        ],
        "img_prompt": "System Prompt编辑器界面，左侧：Prompt编辑区域，显示示例Prompt：\"你是厦门国贸资深供应链风控专家。基于提供的知识库回答问题。规则：1. 只基于上下文回答；2. 不知道明确告知；3. 输出表格形式对比风险点。\"。右侧：提示词块下拉菜单，显示：task_analysis、reflect、sufficiency_check。下方写中文：关键 - 明确人设、严格约束、规范输出。"
    },
    {
        "id": 19, 
        "title": "核心参数：Temperature 配置", 
        "content": [
            "• 低温（0.1-0.3）：严谨、准确、事实查询，企业知识库推荐",
            "• 中温（0.3-0.6）：平衡，一般对话",
            "• 高温（0.7-1.0）：发散、创意，写文案、头脑风暴",
            "• 厦门国贸场景：Temperature < 0.3"
        ],
        "img_prompt": "Temperature仪表盘示意图，从左到右颜色渐变：蓝色（0.1-0.3）→ 青色（0.3-0.6）→ 橙色（0.7-1.0）。指针指向0.2（蓝色区域）。每个区域标注：蓝色区域写中文“严谨准确-事实查询-推荐<0.3”；青色区域写中文“平衡-一般对话”；橙色区域写中文“发散创意-文案写作”。底部红色框写中文：厦门国贸场景 - Temperature < 0.3。"
    },
    {
        "id": 20, 
        "title": "多业务知识库编排", 
        "content": [
            "• 知识库分类：合同库、政策库、产品库、财务库、法务库",
            "• 跨库检索：同时搜索多个知识库，设置优先级",
            "• 权限管控：Private/Team/Public，不同部门不同范围",
            "• 标签体系：Tag Set自动标签，增强检索"
        ],
        "img_prompt": "知识库编排示意图，多个知识库框：合同库（蓝色）、政策库（绿色）、产品库（橙色）、财务库（紫色）、法务库（红色）。上方：检索请求，箭头指向所有知识库，标注“跨库检索”。下方：优先级设置，政策库优先级1（最高）、合同库2、产品库3。右侧：权限设置，Private仅自己、Team团队、Public公开。"
    },
    {
        "id": 21, 
        "title": "数据源同步：Confluence/S3/Notion", 
        "content": [
            "• 支持同步：Confluence、S3、Notion、Discord、Google Drive",
            "• 配置方式：知识库→Configuration→Data Source",
            "• 同步频率：每小时/每天/每周，自动检测变化",
            "• 增量更新：只同步变更文件，节省时间和资源"
        ],
        "img_prompt": "数据源同步架构图，左侧：外部数据源图标：Confluence、Amazon S3、Notion、Discord、Google Drive。中间：RAGFlow同步引擎，写中文“自动同步-增量更新”。右侧：RAGFlow知识库。配置参数显示：同步频率-每天、检测变化-是、增量更新-是。配说明文字：支持10+数据源，自动保持知识库最新。"
    },
    {
        "id": 22, 
        "title": "API 集成：业务系统对接", 
        "content": [
            "• HTTP API：RESTful接口，完整功能覆盖",
            "• Python SDK：pip install ragflow-sdk，简洁易用",
            "• 核心接口：创建知识库、上传文档、检索、对话",
            "• 嵌入场景：OA系统、企业微信、钉钉、业务门户"
        ],
        "img_prompt": "API集成架构图，左侧：业务系统图标：OA系统、企业微信、钉钉、业务门户。中间：API网关，显示两种方式：HTTP API（RESTful）、Python SDK（ragflow-sdk）。右侧：RAGFlow核心服务。下方代码示例：from ragflow_sdk import RAGFlow; client = RAGFlow(api_key='xxx'); kb = client.create_dataset(name='合同库'); kb.upload_document('contract.pdf')。"
    },
    {
        "id": 23, 
        "title": "部署架构：Docker Compose", 
        "content": [
            "• 核心服务：API（Flask）、Web（React）、MySQL、Redis",
            "• 文档引擎：Elasticsearch（默认）或 Infinity（轻量）",
            "• 对象存储：MinIO（本地）或 S3（云端）",
            "• 资源要求：CPU≥4核，RAM≥16GB，Disk≥50GB"
        ],
        "img_prompt": "Docker Compose部署架构图，各服务用容器图标表示。前端：Web（React）；后端：API（Flask/Quart）；数据库：MySQL（元数据）；缓存：Redis；文档引擎：Elasticsearch/Infinity；对象存储：MinIO/S3。用箭头连接显示数据流。底部写资源要求：CPU≥4核，RAM≥16GB，Disk≥50GB，Docker≥24.0。"
    },
    {
        "id": 24, 
        "title": "项目交付验收标准", 
        "content": [
            "• 准确率：Top-5召回率 > 90%，MRR > 0.8",
            "• 响应速度：首字生成 < 3秒，完整答案 < 10秒",
            "• 稳定性：7x24小时高可用，故障率 < 0.1%",
            "• 用户满意度：试用用户满意度 > 4.5/5分"
        ],
        "img_prompt": "验收标准仪表盘，四个指标卡片。第一个卡片：准确率，显示Top-5召回率92%，MRR 0.85，绿色对勾。第二个卡片：响应速度，显示首字生成2.1秒，完整答案7.8秒，绿色对勾。第三个卡片：稳定性，显示7x24可用，故障率0.05%，绿色对勾。第四个卡片：用户满意度，显示4.7/5分，绿色对勾。背景是庆祝图标。"
    },
    {
        "id": 25, 
        "title": "常见问题排障", 
        "content": [
            "• 乱码问题：检查编码，换用UTF-8，启用PaddleOCR",
            "• 表格识别差：切换Table模板，调整table_min_rows=2",
            "• 检索不相关：调整Chunk Size/Overlap，降低相似度阈值，启用Rerank",
            "• 回答幻觉：降低Temperature<0.3，加强System Prompt约束"
        ],
        "img_prompt": "排障指南表，四行两列。第一行：问题-中文乱码，解决方案-检查编码UTF-8，启用PaddleOCR。第二行：问题-表格识别差，解决方案-切换Table模板，table_min_rows=2。第三行：问题-检索不相关，解决方案-调整Chunk Size/Overlap，降低阈值，启用Rerank。第四行：问题-回答幻觉，解决方案-Temperature<0.3，加强Prompt约束。每个解决方案旁配绿色对勾图标。"
    },
    {
        "id": 26, 
        "title": "最佳实践总结", 
        "content": [
            "• 解析：根据文档类型选对应模板（合同用Laws，报表用Table）",
            "• 切片：Chunk Size 300-512，Overlap 10%-15%",
            "• 检索：混合检索+Rerank，关键词0.3+向量0.7",
            "• 生成：Temperature<0.3，严格约束Prompt，强制溯源"
        ],
        "img_prompt": "最佳实践清单，四个要点，每个用大字体标注。第一点：解析策略，写中文“按文档类型选模板 - 合同Laws/报表Table/手册Manual”。第二点：切片参数，写中文“Chunk Size 300-512，Overlap 10%-15%”。第三点：检索配置，写中文“混合检索+Rerank，关键词0.3+向量0.7”。第四点：生成参数，写中文“Temperature<0.3，严格约束Prompt”。每个要点旁配金色星星图标。"
    },
    {
        "id": 27, 
        "title": "核心资产沉淀", 
        "content": [
            "• 知识资产化：让隐性知识显性化，文档变数据资产",
            "• 持续进化：用户反馈→Bad Case分析→优化知识库",
            "• 企业护城河：构建不可复制的供应链行业知识图谱",
            "• 长期价值：降低培训成本，提升决策效率，控制风险"
        ],
        "img_prompt": "知识资产沉淀示意图，从下到上：底层-原始文档（合同、政策、报表）；中层-结构化知识（解析、切片、向量化）；上层-智能应用（问答、Agent、图谱）；顶层-核心资产（企业知识库、知识图谱、最佳实践库）。用向上箭头连接，标注持续进化、用户反馈、Bad Case优化。配文字：构建企业护城河。"
    },
    {
        "id": 28, 
        "title": "结语与Q&A", 
        "content": [
            "• 感谢聆听",
            "• 问答环节 (Q&A)",
            "• 后续支持：文档、API、社区",
            "• 期待与您共同探索 RAG 落地新场景"
        ],
        "img_prompt": "专业结语页，画面中央：厦门国贸大厦轮廓剪影，日出背景。上方：3D发光白色中文字“感谢聆听”。下方：较小文字写“Q&A 问答环节”。底部：联系信息区域，写中文“文档：docs.ragflow.io | GitHub：github.com/infiniflow/ragflow | 社区：Discord”。整体风格温馨专业，象征合作与未来。"
    }
]

# --- FUNCTIONS ---
def generate_images():
    print("Starting SEQUENTIAL Image Generation with OpenRouter...")
    
    # Only generate slide 12
    slides_to_generate = [slide for slide in slides_data if slide['id'] == 12]
    
    for slide in slides_to_generate:
        idx = slide['id']
        filename = f"slide_{idx:02d}.png"
        filepath = os.path.join(IMAGE_PATH, filename)
        
        if os.path.exists(filepath):
            print(f"  Slide {idx} already exists. Skipping.")
            continue
            
        print(f"Generating image {idx}/{len(slides_to_generate)}: {slide['title']}...")
        
        max_retries = 5
        retry_delay = 20 
        
        for attempt in range(max_retries):
            try:
                full_prompt = f"{STYLE_CORE} 重点主体：{slide['img_prompt']} 确保图像包含所有指定的中文字符，文字清晰准确无乱码，流程图、架构图、表格完整可见。高质量3D渲染，16:9比例。"
                
                response = openai_client.chat.completions.create(
                    model="google/gemini-3.1-flash-image-preview",
                    messages=[{"role": "user", "content": full_prompt}],
                    extra_body={"modalities": ["image", "text"]}
                )
                
                result = response.choices[0].message
                success = False
                
                if hasattr(result, 'images') and result.images:
                    for image in result.images:
                        image_url = image['image_url']['url']
                        if image_url.startswith('data:image'):
                            # Extract base64 data
                            header, b64data = image_url.split(',', 1)
                            image_data = base64.b64decode(b64data)
                            with open(filepath, "wb") as f:
                                f.write(image_data)
                            print(f"  [SUCCESS] Slide {idx} saved.")
                            success = True
                            break
                
                if success:
                    break 
                else:
                    print(f"  [WARNING] Attempt {attempt+1} failed for Slide {idx}: No image data.")
                    print(f"    Response: {result}")
            except Exception as e:
                print(f"  [ERROR] Slide {idx} (Attempt {attempt+1}): {e}")
                if "429" in str(e) or "rate_limit" in str(e).lower():
                    time.sleep(retry_delay)
                    retry_delay *= 2
                else:
                    time.sleep(10)
            time.sleep(5)
        time.sleep(15)

def build_pptx():
    print("Assembling PPTX...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    WHITE = RGBColor(255, 255, 255)

    for slide_info in slides_data[3:]:  # From slide 4 onwards
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img_path = os.path.join(IMAGE_PATH, f"slide_{slide_info['id']:02d}.png")
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(5.3), Inches(0), Inches(8.03), Inches(7.5))
        
        bg = slide.shapes.add_shape(1, 0, 0, Inches(5.3), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(5, 12, 25)
        bg.line.width = 0

        t_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(4.5), Inches(1.5))
        t_box.text_frame.word_wrap = True
        p = t_box.text_frame.paragraphs[0]
        p.text = slide_info['title']
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        y = Inches(2.0)
        for bullet in slide_info['content']:
            b_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(4.3), Inches(0.8))
            b_box.text_frame.word_wrap = True
            p = b_box.text_frame.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(200, 200, 200)
            y += Inches(0.8)

        if "subtitle" in slide_info:
            s_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(4.3), Inches(1))
            p = s_box.text_frame.paragraphs[0]
            p.text = slide_info['subtitle']
            p.font.size = Pt(13)
            p.font.italic = True
            p.font.color.rgb = RGBColor(120, 120, 255)

    prs.save(PPT_FILE)
    print(f"Final PPT saved at: {PPT_FILE}")

if __name__ == "__main__":
    generate_images()
    build_pptx()
